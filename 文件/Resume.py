import gspread
from oauth2client.service_account import ServiceAccountCredentials
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import os
import html
import re

# Google Sheets setup
scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
creds_path = r"C:\Henvita\hokjun-5080e816b6a3.json"
creds = ServiceAccountCredentials.from_json_keyfile_name(creds_path, scope)
client = gspread.authorize(creds)

# Open the Google Sheet
sheet_url = "https://docs.google.com/spreadsheets/d/1IaTuqF3NaMfZ_63oL0cdPTFFeEckuZGyftp-oV-rpq4/edit?gid=0#gid=0"
spreadsheet = client.open_by_url(sheet_url)
sheets = spreadsheet.worksheets()  # Get all worksheets
all_data = []  # Store data from all sheets
for sheet in sheets:
    data = sheet.get_all_values()  # Get all data as a list of lists
    all_data.append((sheet.title, data))  # Store sheet title and data


# Generate Word document
def create_word_doc():
    try:
        doc = Document()
        doc.add_heading('Google Sheet Data', 0)
        for sheet_title, data in all_data:
            doc.add_heading(sheet_title, level=1)
            table = doc.add_table(rows=len(data), cols=len(data[0]))
            table.style = 'Table Grid'

            # Add headers
            for j, header in enumerate(data[0]):
                table.cell(0, j).text = header

            # Add data
            for i, row in enumerate(data[1:]):
                for j, cell in enumerate(row):
                    table.cell(i + 1, j).text = cell

        doc.save('sheet_data.docx')
        print("Word document created: sheet_data.docx")
    except Exception as e:
        print(f"Error creating Word document: {e}")


# Generate PowerPoint
def create_ppt():
    try:
        prs = Presentation()
        for sheet_title, data in all_data:
            slide_layout = prs.slide_layouts[5]  # Blank slide with title
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = sheet_title

            # Add table
            rows_count = min(len(data), 10)  # Limit to 10 rows to fit slide
            cols_count = len(data[0])
            table = slide.shapes.add_table(rows_count, cols_count, Inches(1), Inches(1.5), Inches(8), Inches(4)).table

            # Add headers
            for j, header in enumerate(data[0]):
                table.cell(0, j).text = header

            # Add data (up to 10 rows)
            for i in range(1, rows_count):
                for j, cell in enumerate(data[i]):
                    table.cell(i, j).text = cell

        prs.save('sheet_data.pptx')
        print("PowerPoint created: sheet_data.pptx")
    except Exception as e:
        print(f"Error creating PowerPoint: {e}")


# Generate single HTML file with all sheets
def create_html_file():
    # CSS (kept from previous, adjusted for single page)
    css = """
body { font-family: sans-serif; margin: 0; padding: 0; background-color: #000; color: #fff; display: flex; min-height: 100vh; }
* { scrollbar-width: none; -ms-overflow-style: none; }
*::-webkit-scrollbar { display: none; }
.sidebar { width: 133px; background-color: #222; color: #fff; padding: 20px; border-right: 1px solid #444; display: flex; flex-direction: column; position: fixed; top: 0; bottom: 0; overflow-y: auto; }
.sidebar h2 { color: #eee; margin-top: 0; margin-bottom: 15px; }
.sidebar ul { list-style: none; padding: 0; margin: 0; }
.sidebar li a { display: block; padding: 10px 15px; text-decoration: none; color: #ddd; border-radius: 5px; margin-bottom: 5px; cursor: pointer; }
.sidebar li a:hover { background-color: #555; color: #fff; }
.sidebar li a.active { background-color: #777; color: #fff; }
.content { flex-grow: 1; padding: 20px; background-color: #111; margin-left: 173px; }
.table-container { width: 100%; background-color: #333; border-radius: 5px; box-shadow: 0 2px 5px rgba(0, 0, 0, 0.3); padding: 10px; box-sizing: border-box; margin-bottom: 20px; display: none; }
.table-container.active { display: block; }
.table-container table { width: 100%; border-collapse: collapse; color: #eee; }
.table-container th, .table-container td { border: 1px solid #555; padding: 8px; text-align: left; }
.table-container th { background-color: #444; }
@media (max-width: 768px) {
    .sidebar { width: 100px; }
    .content { margin-left: 140px; }
}
"""

    # JavaScript for smooth scrolling and active link highlighting
    js = """
document.addEventListener('DOMContentLoaded', function() {
    const links = document.querySelectorAll('.sidebar li a');
    const containers = document.querySelectorAll('.table-container');

    links.forEach(link => {
        link.addEventListener('click', function(e) {
            e.preventDefault();
            const targetId = this.getAttribute('href').substring(1);
            const targetElement = document.getElementById(targetId);
            if (targetElement) {
                containers.forEach(container => container.classList.remove('active'));
                targetElement.classList.add('active');
                links.forEach(l => l.classList.remove('active'));
                this.classList.add('active');
            }
        });
    });

    // Set first sheet as active on load
    if (links.length > 0 && containers.length > 0) {
        links[0].classList.add('active');
        containers[0].classList.add('active');
    }
});
"""

    # Debug: Print sheet titles to verify
    print("Sheet titles in all_data:", [title for title, _ in all_data])

    def sanitize_id(title, index):
        # Add index to ensure unique IDs
        return f"{re.sub(r'[^a-z0-9-]', '-', title.lower().strip())}-{index}"

    # Sidebar HTML with in-page anchor links
    sidebar_items = "\n".join(
        [f'<li><a href="#{sanitize_id(title, i)}" {"class=\"active\"" if i == 0 else ""}>{html.escape(title)}</a></li>'
         for i, (title, _) in enumerate(all_data)]
    )

    # Table HTML for all sheets
    table_containers = "\n".join(
        [f'''
<div id="{sanitize_id(title, i)}" class="table-container{" active" if i == 0 else ""}">
    <h2>{html.escape(title)}</h2>
    <table>
        <thead><tr>{"".join(f'<th>{html.escape(header)}</th>' for header in data[0])}</tr></thead>
        <tbody>{"".join(f'<tr>{"".join(f'<td>{html.escape(cell)}</td>' for cell in row)}</tr>' for row in data[1:])}</tbody>
    </table>
</div>
''' for i, (title, data) in enumerate(all_data)]
    )

    # Full HTML content
    html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <title>Google Sheet Data</title>
    <style>{css}</style>
    <script>{js}</script>
</head>
<body>
    <div class="sidebar">
        <h2>简历</h2>
        <ul>{sidebar_items}</ul>
    </div>
    <div class="content">
        {table_containers}
    </div>
</body>
</html>
"""

    # Save to file
    try:
        with open('sheet_data.html', 'w', encoding='utf-8') as f:
            f.write(html_content)
        print("HTML file created: sheet_data.html")
    except Exception as e:
        print(f"Error creating HTML file: {e}")


if __name__ == "__main__":
    create_word_doc()
    create_ppt()
    create_html_file()

